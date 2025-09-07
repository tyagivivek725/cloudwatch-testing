#!/usr/bin/env python3
# -*- coding: utf-8 -*-

"""
CloudWatch Alarms → PPT (one slide per alarm)
- Uses get-metric-data (fallback to get-metric-statistics)
- Fixed picture/table rectangles to match your template (cm → inches)
- Branded chart styling (deep blue line, soft fill, dashed red threshold)
- Slide-1 date: replaces [DATE] or any date-looking text
"""

import os
import re
import sys
import json
import tempfile
import subprocess
from datetime import datetime, timedelta, timezone
from pathlib import Path

# Matplotlib (headless)
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.dates as mdates

# python-pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

VERSION = "report_ppt v5.0 (fixed rectangles + styled charts)"

# --- Env / Args ---
REGION       = os.environ.get("AWS_REGION", "ap-south-1")
PROFILE      = os.environ.get("AWS_PROFILE", "default")
JSON_PATH    = sys.argv[1]
PPT_PATH     = sys.argv[2]
WINDOW_HOURS = int(os.environ.get("CW_WINDOW_HOURS", "24"))
TEMPLATE_PATH = os.environ.get("PPT_TEMPLATE")  # optional

# --- Helpers ---
def run_cli(args):
    return subprocess.check_output(args, text=True)

def read_alarms(p):
    with open(p, "r", encoding="utf-8") as f:
        return json.load(f)

def extract_stat_fields(alarm):
    period = alarm.get("Period") or 300
    stat = alarm.get("Statistic") or alarm.get("ExtendedStatistic") or "Average"
    if stat not in ["SampleCount", "Average", "Sum", "Minimum", "Maximum"]:
        stat = "Average"
    return period, stat

def fmt_dims(dims):
    return "; ".join([f"{d['Name']}={d['Value']}" for d in (dims or [])])

def to_dt(ts):
    try:
        return datetime.fromisoformat(ts)
    except Exception:
        # fallback for "YYYY-MM-DDTHH:MM:SSZ"
        return datetime.strptime(ts.replace("Z", ""), "%Y-%m-%dT%H:%M:%S")

# Prefer get-metric-data; fallback to get-metric-statistics
def fetch_series(namespace, metric, dims, period, stat):
    end = datetime.now(timezone.utc)
    start = end - timedelta(hours=WINDOW_HOURS)

    # 1) get-metric-data
    mdq = [{
        "Id": "m1",
        "MetricStat": {
            "Metric": {
                "Namespace": namespace,
                "MetricName": metric,
                "Dimensions": dims or []
            },
            "Period": int(period),
            "Stat": str(stat)
        },
        "ReturnData": True
    }]
    try:
        cmd = [
            "aws", "cloudwatch", "get-metric-data",
            "--metric-data-queries", json.dumps(mdq),
            "--start-time", start.isoformat(),
            "--end-time", end.isoformat(),
            "--scan-by", "TimestampAscending",
            "--region", REGION, "--profile", PROFILE
        ]
        data = json.loads(run_cli(cmd))
        results = (data.get("MetricDataResults") or [])
        if results and results[0].get("Timestamps"):
            pairs = sorted(
                zip(results[0]["Timestamps"], results[0].get("Values", [])),
                key=lambda t: t[0]
            )
            xs = [to_dt(t.isoformat() if hasattr(t, "isoformat") else str(t)) for t, _ in pairs]
            ys = [v for _, v in pairs]
            return xs, ys
    except Exception:
        pass

    # 2) fallback: get-metric-statistics
    try:
        base = [
            "aws", "cloudwatch", "get-metric-statistics",
            "--namespace", namespace,
            "--metric-name", metric,
            "--start-time", start.isoformat(),
            "--end-time", end.isoformat(),
            "--period", str(period),
            "--statistics", str(stat),
            "--region", REGION, "--profile", PROFILE
        ]
        for d in (dims or []):
            base += ["--dimensions", f"Name={d['Name']},Value={d['Value']}"]
        data = json.loads(run_cli(base))
        pts = sorted(data.get("Datapoints", []), key=lambda x: x["Timestamp"])
        xs = [to_dt(p["Timestamp"]) for p in pts]
        ys = [p.get(stat, p.get("Average", 0.0)) for p in pts]
        return xs, ys
    except Exception:
        pass

    return [], []

# --- Chart styling ---
def plot_chart(xs, ys, ylabel, threshold=None, out_png=None, note=None):
    """
    Branded, high-contrast chart:
    - 12.20in x 3.02in @ 300 DPI (matches 30.97 x 7.66 cm picture rectangle)
    - Deep blue line, soft fill, red dashed threshold
    - Last-point marker + label
    """
    out_png = out_png or Path(tempfile.mkdtemp()) / "chart.png"

    # Palette
    LINE_COLOR = "#003087"   # deep blue
    FILL_COLOR = "#003087"   # same hue, transparent fill
    GRID_COLOR = "#B0BEC5"   # light gray
    THR_COLOR  = "#D0021B"   # red for thresholds
    TEXT_GRAY  = "#455A64"

    fig = plt.figure(figsize=(12.20, 3.02), dpi=300)
    ax = fig.add_subplot(111)

    if xs and ys:
        # main series
        ax.plot(xs, ys, linewidth=2.8, color=LINE_COLOR, antialiased=True)
        # soft fill under the line
        try:
            bottom = min(ys) if ys else 0.0
            ax.fill_between(xs, ys, [bottom] * len(ys), alpha=0.15, color=FILL_COLOR)
        except Exception:
            pass

        # last-point marker + label
        try:
            ax.plot(xs[-1], ys[-1], "o", ms=6, color=LINE_COLOR)
            ax.annotate(
                f"{ys[-1]:.2f}",
                xy=(xs[-1], ys[-1]),
                xytext=(6, -10),
                textcoords="offset points",
                fontsize=10,
                color=TEXT_GRAY,
            )
        except Exception:
            pass
    else:
        # Empty-state visual
        ax.plot([], [])
        ax.text(
            0.5, 0.5, note or "No datapoints in selected window",
            transform=ax.transAxes,
            ha="center", va="center",
            fontsize=13, color=TEXT_GRAY
        )

    # Grid & axes styling
    ax.grid(True, linewidth=0.6, color=GRID_COLOR, alpha=0.7)
    for spine in ax.spines.values():
        spine.set_linewidth(0.8)
        spine.set_color("#90A4AE")  # muted grey

    ax.set_xlabel("Time (last 24h)", fontsize=12)
    ax.set_ylabel(ylabel, fontsize=12)

    ax.tick_params(axis="both", labelsize=10)
    ax.xaxis.set_major_locator(mdates.AutoDateLocator(minticks=6, maxticks=10))
    ax.xaxis.set_major_formatter(mdates.DateFormatter("%d %b %H:%M"))
    fig.autofmt_xdate(rotation=0)

    # Threshold
    if threshold is not None:
        try:
            thr = float(threshold)
            ax.axhline(thr, linewidth=1.8, linestyle="--", color=THR_COLOR)
        except Exception:
            pass

    # breathing room for tick labels
    fig.subplots_adjust(left=0.08, right=0.995, top=0.98, bottom=0.24)
    fig.savefig(out_png, dpi=300, bbox_inches=None, pad_inches=0.05)
    plt.close(fig)
    return str(out_png)

# --- Slide-1 date updater (replaces [DATE] or any date-looking text) ---
MONTHS = (
    "January", "February", "March", "April", "May", "June",
    "July", "August", "September", "October", "November", "December"
)
DATE_RX = re.compile(
    r"\b(\d{1,2})(st|nd|rd|th)[-\s]"
    r"(January|February|March|April|May|June|July|August|September|October|November|December)"
    r"[-\s](\d{4})\b",
    re.IGNORECASE
)

def ordinal(n: int) -> str:
    return "th" if 10 <= n % 100 <= 20 else {1: "st", 2: "nd", 3: "rd"}.get(n % 10, "th")

def today_str() -> str:
    now = datetime.now().astimezone()
    return f"{now.day:02d}{ordinal(now.day)}-{MONTHS[now.month-1]}-{now.year}"

def update_first_slide_date(prs: Presentation) -> None:
    if not prs.slides:
        return
    s = prs.slides[0]
    new_text = today_str()
    for shape in s.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            txt = shape.text_frame.text
        except Exception:
            continue
        if "[DATE]" in txt:
            shape.text_frame.clear()
            shape.text_frame.paragraphs[0].text = txt.replace("[DATE]", new_text)
            return
        if DATE_RX.search(txt):
            shape.text_frame.clear()
            shape.text_frame.paragraphs[0].text = DATE_RX.sub(new_text, txt)
            return

# --- Layout helpers ---
def first_layout_with_title(prs: Presentation):
    for layout in prs.slide_layouts:
        types = set(
            sh.placeholder_format.type
            for sh in layout.placeholders
            if sh.is_placeholder
        )
        if (PP_PLACEHOLDER.TITLE in types) or (PP_PLACEHOLDER.CENTER_TITLE in types):
            return layout
    return prs.slide_layouts[0]

def remove_non_title_placeholders(slide):
    to_remove = []
    for shp in slide.shapes:
        try:
            if shp.is_placeholder:
                t = shp.placeholder_format.type
                if t not in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                    to_remove.append(shp)
        except Exception:
            continue
    for shp in to_remove:
        el = shp._element
        el.getparent().remove(el)

# --- Build deck ---
def build_ppt(alarms, out_path):
    # Load template (if provided), else blank 16:9
    if TEMPLATE_PATH and os.path.isfile(TEMPLATE_PATH):
        prs = Presentation(TEMPLATE_PATH)
    else:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    update_first_slide_date(prs)

    header_fill = RGBColor(0x18, 0x4A, 0x5A)
    zebra = RGBColor(0xF0, 0xF3, 0xF5)
    content_layout = first_layout_with_title(prs)

    for a in alarms:
        ns = a.get("Namespace")
        metric = a.get("MetricName")
        dims = a.get("Dimensions", [])
        title = a.get("AlarmName", "(no name)")
        threshold = a.get("Threshold")
        period, stat = extract_stat_fields(a)

        slide = prs.slides.add_slide(content_layout)
        remove_non_title_placeholders(slide)

        if slide.shapes.title:
            slide.shapes.title.text = title
            try:
                slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(26)
            except Exception:
                pass

        # ----- EXACT POSITIONS (cm → inches) -----
        # Chart (Format Picture): Height 7.66 cm, Width 30.97 cm, X 1.04 cm, Y 2.55 cm
        chart_left   = Inches(1.04 / 2.54)
        chart_top    = Inches(2.55 / 2.54)
        chart_width  = Inches(30.97 / 2.54)
        chart_height = Inches(7.66 / 2.54)

        # Table (Format Shape):  Height 7.62 cm, Width 30.15 cm, X 1.86 cm, Y 10.21 cm
        tbl_left   = Inches(1.86 / 2.54)
        tbl_top    = Inches(10.21 / 2.54)
        tbl_width  = Inches(30.15 / 2.54)
        tbl_height = Inches(7.62 / 2.54)
        # ----------------------------------------

        # Fetch series
        if not (ns and metric):
            xs, ys = [], []
            note = "Composite alarm / no base metric"
        else:
            xs, ys = fetch_series(ns, metric, dims, period, stat)
            note = None if xs else f"No datapoints in last {WINDOW_HOURS}h ({REGION})"

        # Place chart
        chart_png = plot_chart(xs, ys, ylabel=metric or "", threshold=threshold, note=note)
        slide.shapes.add_picture(
            chart_png, chart_left, chart_top, width=chart_width, height=chart_height
        )

        # Build styled details table
        rows = [
            ("Metric",      metric or "composite/none"),
            ("Namespace",   ns or ""),
            ("Statistic",   a.get("Statistic", "") or a.get("ExtendedStatistic", "Average")),
            ("Threshold",   threshold),
            ("Operator",    a.get("ComparisonOperator", "")),
            ("EvalPeriods", a.get("EvaluationPeriods", "")),
            ("Dimensions",  fmt_dims(dims)),
            ("Updated",     a.get("StateUpdatedTimestamp", "")),
        ]
        n_rows = len(rows) + 1

        table_shape = slide.shapes.add_table(
            rows=n_rows, cols=2, left=tbl_left, top=tbl_top,
            width=tbl_width, height=tbl_height
        )
        table = table_shape.table
        table.columns[0].width = int(tbl_width * 0.30)
        table.columns[1].width = int(tbl_width * 0.70)

        # Header row
        table.cell(0, 0).text = "Metric Detail"
        table.cell(0, 1).text = "Value"
        for c in range(2):
            cell = table.cell(0, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill
            tf = cell.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = "Metric Detail" if c == 0 else "Value"
            p.font.bold = True
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            p.alignment = PP_ALIGN.LEFT

        # Body rows (zebra)
        for i, (label, value) in enumerate(rows, start=1):
            lc = table.cell(i, 0)
            vc = table.cell(i, 1)
            if i % 2 == 1:
                lc.fill.solid(); lc.fill.fore_color.rgb = zebra
                vc.fill.solid(); vc.fill.fore_color.rgb = zebra
            lc.text = str(label)
            lp = lc.text_frame.paragraphs[0]
            lp.font.bold = True
            lp.font.size = Pt(12)
            vc.text = "" if value is None else str(value)
            tf = vc.text_frame
            tf.word_wrap = True
            for para in tf.paragraphs:
                para.font.size = Pt(12)
                para.alignment = PP_ALIGN.LEFT

    prs.save(PPT_PATH)

# --- Main ---
def main():
    print(VERSION)
    data = read_alarms(JSON_PATH)
    metric_alarms = data.get("MetricAlarms", []) or []
    composite_alarms = data.get("CompositeAlarms", []) or []
    all_alarms = metric_alarms + composite_alarms
    build_ppt(all_alarms, PPT_PATH)
    print(f"Wrote PPT: {PPT_PATH}")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: report_ppt.py <alarms.json> <out.pptx>")
        sys.exit(2)
    main()
